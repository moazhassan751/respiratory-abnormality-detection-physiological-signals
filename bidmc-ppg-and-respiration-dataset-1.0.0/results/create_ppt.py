"""
Create PowerPoint Presentation for First Evaluation
Respiratory Abnormality Detection using PPG Signals
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
import os

# Alias for convenience
RgbColor = RGBColor

def add_cover_slide(prs, department, course, instructor, members):
    """Add a cover slide with department and group info"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Department
    dept_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12.333), Inches(0.8))
    tf = dept_box.text_frame
    p = tf.paragraphs[0]
    p.text = department
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Course
    course_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(0.6))
    tf = course_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Course: {course}"
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Instructor
    inst_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(12.333), Inches(0.6))
    tf = inst_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Instructor: {instructor}"
    p.font.size = Pt(22)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    # Divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(3.6), Inches(7.333), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LIGHT_BLUE
    line.line.fill.background()
    
    # Group Members Title
    grp_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.6))
    tf = grp_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "GROUP MEMBERS"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Members
    members_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(12.333), Inches(2))
    tf = members_box.text_frame
    tf.word_wrap = True
    for i, member in enumerate(members):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = member
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)
    
    return slide

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BLUE = RgbColor(26, 82, 118)      # #1a5276
LIGHT_BLUE = RgbColor(52, 152, 219)    # #3498db
WHITE = RgbColor(255, 255, 255)
DARK_GRAY = RgbColor(51, 51, 51)
GREEN = RgbColor(39, 174, 96)          # #27ae60

def add_title_slide(prs, title, subtitle):
    """Add a title slide with gradient-like background"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def add_table_slide(prs, title, headers, rows):
    """Add a slide with a table"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    
    col_width = Inches(12.333 / num_cols)
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)).table
    
    # Set column widths
    for i in range(num_cols):
        table.columns[i].width = col_width
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_BLUE
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
    
    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RgbColor(248, 249, 250)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a two-column slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    title_bar.fill.solid()
    title_bar.fill.fore_color.rgb = DARK_BLUE
    title_bar.line.fill.background()
    
    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Left column content
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.9), Inches(6), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.4), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    # Right column content
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.9), Inches(6), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

# ============================================================================
# SLIDE 1: Cover Slide (Department & Group Info)
# ============================================================================
add_cover_slide(
    prs,
    "DEPARTMENT: CREATIVE TECHNOLOGIES",
    "Biomedical Data Processing and Analysis",
    "Dr. Abdul Haleem Butt",
    [
        "MOAZ HASSAN (231168)",
        "HUNAIN AHMED (231156)",
        "MAHAD JOKHIO (231242)"
    ]
)

# ============================================================================
# SLIDE 2: Title Slide
# ============================================================================
add_title_slide(
    prs,
    "RESPIRATORY ABNORMALITY DETECTION\nUSING PHYSIOLOGICAL SIGNAL PROCESSING",
    "A Machine Learning Approach for ICU Patient Monitoring\n\nFirst Evaluation: Data Acquisition to Feature Selection\nDecember 2025"
)

# ============================================================================
# SLIDE 2: Problem Statement & Objectives
# ============================================================================
add_two_column_slide(
    prs,
    "Problem Statement & Objectives",
    "The Problem",
    [
        "Invasive respiratory monitoring causes patient discomfort",
        "Expensive equipment limits accessibility",
        "Hospital-only monitoring, no continuous home monitoring",
        "480 million affected globally need early detection"
    ],
    "Our Objectives",
    [
        "✓ Acquire and analyze BIDMC PPG dataset",
        "✓ Implement state-of-the-art preprocessing",
        "✓ Extract comprehensive multi-domain features",
        "✓ Select optimal features for classification",
        "→ Build and evaluate classification model (Next Phase)"
    ]
)

# ============================================================================
# SLIDE 3: Literature Review & Research Gap
# ============================================================================
add_table_slide(
    prs,
    "Literature Review & Research Gap",
    ["Study", "Year", "Method", "Accuracy", "Limitation"],
    [
        ["Charlton et al.", "2016", "Frequency analysis", "85%", "Limited features"],
        ["Pimentel et al.", "2017", "Time-domain only", "78%", "No wavelet analysis"],
        ["Birrenkott et al.", "2018", "Deep Learning", "88%", "No interpretability"],
        ["Liu et al.", "2020", "Hybrid approach", "90%", "Small dataset (n=20)"],
        ["", "", "", "", ""],
        ["OUR APPROACH", "", "Multi-domain", "Target", "101 features, 7-step"],
        ["", "", "features + RF", ">90%", "SOTA preprocessing"]
    ]
)

# ============================================================================
# SLIDE 4: Dataset Description
# ============================================================================
add_table_slide(
    prs,
    "Dataset: BIDMC PPG and Respiration (PhysioNet)",
    ["Attribute", "Value", "Signal", "Use in Project"],
    [
        ["Source", "PhysioNet", "PPG", "Primary input signal"],
        ["Subjects", "53 ICU patients", "ECG", "Heart rate validation"],
        ["Duration", "8 min/subject", "RESP", "Ground truth reference"],
        ["Sampling Rate", "125 Hz", "SpO2/HR", "Clinical features"],
        ["Total Samples", "~3.18 million", "", ""],
    ]
)

# ============================================================================
# SLIDE 5: Preprocessing Pipeline
# ============================================================================
add_table_slide(
    prs,
    "Preprocessing Pipeline (7 Steps - SOTA)",
    ["Step", "Method", "Purpose", "Result"],
    [
        ["1. Signal Quality", "Correlation SQI", "Remove unreliable segments", ""],
        ["2. Missing Values", "Linear interpolation", "Handle data gaps", ""],
        ["3. Baseline Removal", "Median filter (0.5Hz)", "Remove DC drift", "SNR: +133%"],
        ["4. Powerline Removal", "Notch filter (50/60Hz)", "Remove electrical noise", "Artifacts: -94%"],
        ["5. Bandpass Filter", "Butterworth (0.1-8Hz)", "Preserve respiratory band", "Missing: -100%"],
        ["6. Artifact Removal", "Hampel filter (MAD)", "Remove motion artifacts", ""],
        ["7. Normalization", "Z-score", "Comparable across subjects", ""],
    ]
)

# ============================================================================
# SLIDE 6: Feature Extraction Overview
# ============================================================================
add_table_slide(
    prs,
    "Feature Extraction: 101 Multi-Domain Features",
    ["Domain", "# Features", "Key Features"],
    [
        ["Time Domain", "15", "Mean, Std, Skewness, Kurtosis, RMS"],
        ["Frequency Domain", "12", "Dominant freq, LF/HF power, Spectral entropy"],
        ["Wavelet Domain", "20", "Daubechies-4 (5 levels): A5, D5 energy"],
        ["HRV Features", "18", "SDNN, RMSSD, pNN50, LF/HF ratio"],
        ["ECG Features", "15", "R-peak variability, QRS statistics"],
        ["Numeric Features", "8", "SpO2 mean/min, Heart rate, Perfusion"],
        ["Ground Truth", "8", "Reference RR, Breath amplitude"],
        ["Demographics", "5", "Age, Duration, Recording info"],
    ]
)

# ============================================================================
# SLIDE 7: Key Feature Extraction Methods
# ============================================================================
add_two_column_slide(
    prs,
    "Key Feature Extraction Methods",
    "Frequency Domain (FFT)",
    [
        "LF Band (0.04-0.15 Hz): Sympathetic activity",
        "HF Band (0.15-0.4 Hz): Parasympathetic (respiratory)",
        "Respiratory Band (0.1-0.5 Hz): Direct modulation",
        "Spectral Entropy: Signal complexity"
    ],
    "Wavelet Domain (db4, 5 levels)",
    [
        "D5 (1.95-3.9 Hz): Cardiac fundamental",
        "A5 (0-1.95 Hz): Respiratory modulation ★",
        "Energy, Entropy per level",
        "RSA: Heart rate varies with breathing"
    ]
)

# ============================================================================
# SLIDE 8: Feature Selection Method
# ============================================================================
add_two_column_slide(
    prs,
    "Feature Selection: Random Forest Importance",
    "Why Feature Selection?",
    [
        "Curse of dimensionality with 101 features",
        "Overfitting risk with small dataset",
        "Computational cost reduction needed",
        "Clinical interpretability required"
    ],
    "Random Forest Advantages",
    [
        "Handles non-linear relationships",
        "Built-in feature importance scores",
        "Robust to outliers (tree-based)",
        "Manages correlated features well",
        "n_estimators = 100 trees"
    ]
)

# ============================================================================
# SLIDE 9: Feature Selection Results
# ============================================================================
add_table_slide(
    prs,
    "Feature Selection Results: Top 10 Features",
    ["Rank", "Feature", "Importance", "Domain", "Why Selected"],
    [
        ["1", "resp_rate_mean", "0.142", "Ground Truth", "Direct respiratory measure"],
        ["2", "spo2_mean", "0.098", "Numeric", "Oxygen efficiency"],
        ["3", "hrv_hf_power", "0.087", "HRV", "RSA (respiratory linked)"],
        ["4", "wavelet_a5_energy", "0.076", "Wavelet", "Respiratory modulation"],
        ["5", "resp_amplitude_mean", "0.065", "Ground Truth", "Breathing depth"],
        ["6", "ppg_respiratory_freq", "0.058", "Frequency", "Primary resp rate"],
        ["7", "hrv_rmssd", "0.052", "HRV", "Parasympathetic activity"],
        ["8", "breath_duration_std", "0.048", "Time", "Breathing regularity"],
    ]
)

# ============================================================================
# SLIDE 10: Summary & Achievements
# ============================================================================
add_table_slide(
    prs,
    "Summary & Achievements",
    ["Milestone", "Status", "Key Result"],
    [
        ["Data Acquisition", "✓ Complete", "53 subjects, 3.18M samples"],
        ["Preprocessing", "✓ Complete", "7-step SOTA, 133% SNR improvement"],
        ["Feature Extraction", "✓ Complete", "101 features from 6 domains"],
        ["Feature Selection", "✓ Complete", "Top 40 features, 60% reduction"],
        ["", "", ""],
        ["Tools Used", "", "Python, NumPy, Pandas, SciPy, Scikit-learn"],
        ["Output Files", "", "features.csv, feature_importance.csv, clinical_report.txt"],
    ]
)

# ============================================================================
# SLIDE 11: Future Work
# ============================================================================
add_content_slide(
    prs,
    "Future Work (Next Phase)",
    [
        "• Model Building: Random Forest, SVM, XGBoost classifiers",
        "• Cross-Validation: 5-fold stratified CV",
        "• Performance Metrics: Accuracy, Precision, Recall, F1-Score, ROC-AUC",
        "• Hyperparameter Tuning: Grid search optimization",
        "• Clinical Validation: Interpret results for practical use",
        "",
        "Long-term Goals:",
        "• Real-time respiratory monitoring system",
        "• Integration with wearable devices (smartwatches)",
        "• Deployment as mobile/web application"
    ]
)

# ============================================================================
# SLIDE 12: References & Thank You
# ============================================================================
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Background
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

# Thank You text
thank_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(1))
tf = thank_box.text_frame
p = tf.paragraphs[0]
p.text = "THANK YOU!"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.alignment = PP_ALIGN.CENTER

# Questions
q_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(0.8))
tf = q_box.text_frame
p = tf.paragraphs[0]
p.text = "Questions?"
p.font.size = Pt(32)
p.font.color.rgb = RgbColor(200, 200, 200)
p.alignment = PP_ALIGN.CENTER

# References
ref_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(3))
tf = ref_box.text_frame
tf.word_wrap = True
refs = [
    "Key References:",
    "1. Pimentel et al. (2017) - IEEE TBME - Respiratory rate estimation",
    "2. Charlton et al. (2016) - Physiological Measurement - Algorithm assessment",
    "3. Goldberger et al. (2000) - PhysioNet database",
    "4. NeuroKit2 Documentation (2023) - Biosignal processing"
]
for i, ref in enumerate(refs):
    if i == 0:
        p = tf.paragraphs[0]
        p.font.bold = True
    else:
        p = tf.add_paragraph()
    p.text = ref
    p.font.size = Pt(16)
    p.font.color.rgb = RgbColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

# Save presentation
output_path = r"d:\Downloads\bidmc-ppg-and-respiration-dataset-1.0.0-20251011T094008Z-1-001\bidmc-ppg-and-respiration-dataset-1.0.0\results\First_Evaluation_Presentation.pptx"
prs.save(output_path)

print("=" * 60)
print("✅ PowerPoint Presentation Created Successfully!")
print("=" * 60)
print(f"\nFile: {output_path}")
print(f"\nSlides: 13")
print("\nSlide Contents:")
print("  1. Cover Slide (Department & Group Members)")
print("  2. Title Slide")
print("  3. Problem Statement & Objectives")
print("  4. Literature Review & Research Gap")
print("  5. Dataset Description")
print("  6. Preprocessing Pipeline (7 Steps)")
print("  7. Feature Extraction Overview")
print("  8. Key Feature Extraction Methods")
print("  9. Feature Selection Method")
print(" 10. Feature Selection Results")
print(" 11. Summary & Achievements")
print(" 12. Future Work")
print(" 13. References & Thank You")
